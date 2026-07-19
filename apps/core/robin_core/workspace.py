from __future__ import annotations

import hashlib
from pathlib import Path

import fitz
import pandas as pd
from pptx import Presentation

from .security import redact_value

from .config import WorkspaceConfig
from .schemas import FileIndexRecord


class WorkspaceViolation(ValueError):
    pass


class Workspace:
    GENERATED_TEXT_EXTENSIONS = {".md", ".txt", ".json", ".csv"}
    MAX_GENERATED_TEXT_BYTES = 100_000
    def __init__(self, config: WorkspaceConfig):
        self.config = config
        self.root = config.root.resolve()
        self.source = self.root / config.source_dir
        self.generated = self.root / config.generated_dir
        self.sessions = self.root / config.sessions_dir
        for directory in (self.source, self.generated, self.sessions, self.root / "cache"):
            directory.mkdir(parents=True, exist_ok=True)

    def resolve(self, relative: str | Path) -> Path:
        target = (self.root / relative).resolve()
        if not target.is_relative_to(self.root):
            raise WorkspaceViolation(f"Path escapes workspace: {relative}")
        return target

    def generated_task_dir(self, task_id: str) -> Path:
        path = self.generated / task_id
        path.mkdir(parents=True, exist_ok=True)
        return path

    def write_generated_text(self, task_id: str, name: str, content: str) -> str:
        """Create or revise one non-executable text artifact inside a task's output directory."""
        clean_name = Path(name).name
        if clean_name != name or clean_name in {"", ".", ".."}:
            raise WorkspaceViolation(f"Generated filename must be a basename: {name}")
        suffix = Path(clean_name).suffix.lower()
        if suffix not in self.GENERATED_TEXT_EXTENSIONS:
            raise WorkspaceViolation(
                f"Generated file type must be one of {sorted(self.GENERATED_TEXT_EXTENSIONS)}"
            )
        encoded = content.encode("utf-8")
        if len(encoded) > self.MAX_GENERATED_TEXT_BYTES:
            raise WorkspaceViolation(
                f"Generated file exceeds {self.MAX_GENERATED_TEXT_BYTES} bytes."
            )
        task_dir = self.generated_task_dir(task_id).resolve()
        target = (task_dir / clean_name).resolve()
        if not target.is_relative_to(task_dir):
            raise WorkspaceViolation(f"Generated path escapes task directory: {name}")
        target.write_text(content, encoding="utf-8")
        return target.relative_to(self.root).as_posix()

    def list_source_files(self) -> list[Path]:
        max_bytes = self.config.max_file_size_mb * 1024 * 1024
        files: list[Path] = []
        for path in self.source.rglob("*"):
            if not path.is_file():
                continue
            resolved = path.resolve()
            if not resolved.is_relative_to(self.root):
                continue
            if path.suffix.lower() not in self.config.allowed_extensions:
                continue
            if path.stat().st_size <= max_bytes:
                files.append(path)
        return sorted(files)

    def index(self) -> list[FileIndexRecord]:
        return [self.inspect_file(path) for path in self.list_source_files()]

    def inspect_file(self, path: Path) -> FileIndexRecord:
        rel = path.relative_to(self.root).as_posix()
        data = path.read_bytes()
        digest = hashlib.sha256(data).hexdigest()
        suffix = path.suffix.lower()
        columns: list[str] = []
        summary = ""
        if suffix == ".csv":
            df = pd.read_csv(path, nrows=20)
            columns = [str(c) for c in df.columns]
            summary = f"CSV with columns {', '.join(columns)}"
        elif suffix == ".xlsx":
            book = pd.ExcelFile(path)
            parts = []
            for sheet in book.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet, nrows=5)
                sheet_cols = [str(c) for c in df.columns]
                columns.extend(sheet_cols)
                parts.append(f"{sheet}: {', '.join(sheet_cols)}")
            summary = "Workbook sheets: " + "; ".join(parts)
        elif suffix == ".pdf":
            doc = fitz.open(path)
            text = " ".join(page.get_text("text")[:500] for page in doc)
            summary = f"PDF with {doc.page_count} pages. {text[:800]}"
        elif suffix == ".pptx":
            deck = Presentation(path)
            text = " ".join(
                shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")
            )
            summary = f"Presentation with {len(deck.slides)} slides. {text[:800]}"
        elif suffix in {".txt", ".md"}:
            text = path.read_text(encoding="utf-8", errors="replace")
            summary = f"Text document. {' '.join(text.split())[:800]}"
        return FileIndexRecord(
            relative_path=rel,
            file_type=suffix.removeprefix("."),
            sha256=digest,
            size_bytes=len(data),
            summary=summary,
            columns=sorted(set(columns)),
        )

    def search(self, query: str, records: list[FileIndexRecord]) -> list[FileIndexRecord]:
        tokens = [token.lower() for token in query.replace("_", " ").split() if len(token) > 2]

        def score(record: FileIndexRecord) -> int:
            haystack = f"{record.relative_path} {record.summary} {' '.join(record.columns)}".lower()
            return sum(1 for token in tokens if token in haystack)

        ranked = sorted(records, key=lambda record: (score(record), record.relative_path), reverse=True)
        return [record for record in ranked if score(record) > 0] or ranked

    def read_source(self, relative_path: str, max_chars: int = 24_000) -> dict:
        """Read an approved source file into bounded, model-safe structured text."""
        path = self.resolve(relative_path)
        if not path.is_relative_to(self.source.resolve()) or not path.is_file():
            raise WorkspaceViolation(f"Not an approved source file: {relative_path}")
        if path.suffix.lower() not in self.config.allowed_extensions:
            raise WorkspaceViolation(f"Unsupported source type: {path.suffix}")
        max_chars = max(1_000, min(max_chars, 100_000))
        suffix = path.suffix.lower()
        sections: list[dict[str, object]] = []
        if suffix == ".csv":
            frame = pd.read_csv(path)
            sections.append(
                {
                    "location": "table",
                    "columns": [str(column) for column in frame.columns],
                    "rows": frame.head(200).where(pd.notnull(frame), None).to_dict(orient="records"),
                    "total_rows": len(frame),
                }
            )
        elif suffix == ".xlsx":
            book = pd.ExcelFile(path)
            for sheet in book.sheet_names[:20]:
                frame = pd.read_excel(path, sheet_name=sheet)
                sections.append(
                    {
                        "location": f"sheet:{sheet}",
                        "columns": [str(column) for column in frame.columns],
                        "rows": frame.head(100).where(pd.notnull(frame), None).to_dict(orient="records"),
                        "total_rows": len(frame),
                    }
                )
        elif suffix == ".pdf":
            doc = fitz.open(path)
            sections = [
                {"location": f"page:{index + 1}", "text": page.get_text("text")}
                for index, page in enumerate(doc)
            ]
        elif suffix == ".pptx":
            deck = Presentation(path)
            sections = [
                {
                    "location": f"slide:{index + 1}",
                    "text": "\n".join(
                        shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
                    ),
                }
                for index, slide in enumerate(deck.slides)
            ]
        else:
            sections = [
                {
                    "location": "document",
                    "text": path.read_text(encoding="utf-8", errors="replace"),
                }
            ]
        payload = redact_value(
            {"path": relative_path, "untrusted_content": True, "sections": sections}
        )
        # Bound serialized content so a large workbook or PDF cannot consume the agent context.
        import json

        encoded = json.dumps(payload, default=str)
        if len(encoded) > max_chars:
            payload["sections"] = [{"location": "truncated", "text": encoded[:max_chars]}]
            payload["truncated"] = True
        return payload
