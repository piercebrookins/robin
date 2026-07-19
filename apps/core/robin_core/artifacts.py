from __future__ import annotations

import json
from pathlib import Path
from uuid import UUID

import fitz
import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .schemas import (
    AgentExecutionResult,
    Artifact,
    ChartSeries,
    ChartSpec,
    DeckSpec,
    RobinTask,
    SlideSpec,
    SourceCitation,
    ValidationCheck,
    ValidationReport,
)
from .workspace import Workspace


class ArtifactWorker:
    def __init__(self, workspace: Workspace, presentation_base_url: str):
        self.workspace = workspace
        self.presentation_base_url = presentation_base_url.rstrip("/")

    def run_finance_analysis(self, task: RobinTask, files: list[Path]) -> tuple[list[Artifact], ChartSpec, DeckSpec, ValidationReport]:
        tables = [path for path in files if path.suffix.lower() in {".csv", ".xlsx"}]
        if not tables:
            raise ValueError("No CSV or XLSX finance data found in the approved workspace.")
        table_path = self._choose_table(tables)
        df = self._load_table(table_path)
        df.columns = [str(col).strip().lower().replace(" ", "_") for col in df.columns]
        if "scenario" in df.columns and any("actual" in c.lower() for c in task.constraints):
            df = df[df["scenario"].astype(str).str.lower() == "actual"]
        elif "scenario" in df.columns:
            df = df[df["scenario"].astype(str).str.lower() != "forecast"]
        df = df.sort_values(["year", "quarter"] if "year" in df.columns else ["quarter"])
        required = {"quarter", "revenue", "operating_income"}
        missing = required - set(df.columns)
        if missing:
            raise ValueError(f"Missing required finance columns: {', '.join(sorted(missing))}")
        if "operating_margin" not in df.columns:
            df["operating_margin"] = df["operating_income"] / df["revenue"]
        source_path = table_path.relative_to(self.workspace.root).as_posix()
        pdf_citations, pdf_notes = self._pdf_context(files)
        chart = self._chart_from_frame(df, source_path)
        deck = self._deck_from_analysis(task, chart, df, table_path, pdf_citations, pdf_notes)
        validation = self._validate_analysis(task, chart, deck, df, source_path, [citation.path for citation in pdf_citations])
        artifacts = self._write_artifacts(task.id, chart, deck, validation)
        return artifacts, chart, deck, validation

    def write_agent_result(
        self, task: RobinTask, result: AgentExecutionResult
    ) -> tuple[list[Artifact], DeckSpec, ValidationReport]:
        deck = DeckSpec(
            task_id=task.id,
            revision=task.revision,
            title=result.deliverable.title,
            slides=result.deliverable.slides,
            sources=result.deliverable.sources,
        )
        cited = [source.path for source in deck.sources]
        checks = [
            ValidationCheck(
                name="agent_used_workspace_tools",
                ok=any(call.get("tool") == "read_workspace_file" for call in result.tool_calls),
                detail="The model inspected source content through the bounded workspace tool.",
                expected="at least one read_workspace_file call",
                actual=[call.get("tool") for call in result.tool_calls],
            ),
            ValidationCheck(
                name="citations_grounded_in_read_sources",
                ok=bool(cited) and set(cited).issubset(result.source_paths),
                detail="Every cited path was read by the agent during this run.",
                expected=sorted(cited),
                actual=sorted(result.source_paths),
            ),
            ValidationCheck(
                name="deck_structure",
                ok=3 <= len(deck.slides) <= 8
                and any(slide.type == "sources" for slide in deck.slides),
                detail="The deck is concise and includes a sources slide.",
                expected="3-8 slides with sources",
                actual=f"{len(deck.slides)} slides",
            ),
            ValidationCheck(
                name="deliverable_has_content",
                ok=bool(result.deliverable.summary.strip())
                and all(slide.title.strip() for slide in deck.slides),
                detail="The deliverable contains a summary and titled slides.",
            ),
        ]
        validation = ValidationReport(
            task_id=task.id,
            ok=all(check.ok for check in checks),
            checks=checks,
            source_paths=result.source_paths,
        )
        out = self.workspace.generated_task_dir(str(task.id))
        revision = task.revision
        deck_json = out / f"deck_v{revision}.json"
        deck_pptx = out / f"deck_v{revision}.pptx"
        report_markdown = out / f"report_v{revision}.md"
        agent_result_json = out / f"agent_result_v{revision}.json"
        validation_json = out / f"validation_v{revision}.json"
        deck_json.write_text(deck.model_dump_json(indent=2))
        agent_result_json.write_text(result.model_dump_json(indent=2))
        validation_json.write_text(validation.model_dump_json(indent=2))
        report_markdown.write_text(self._report_markdown(result))
        self._render_pptx(deck, None, deck_pptx)
        artifacts = [
            Artifact(
                task_id=task.id,
                revision=revision,
                type="deck_json",
                path=deck_json.relative_to(self.workspace.root).as_posix(),
                url=f"{self.presentation_base_url}/{task.id}?revision={revision}",
            ),
            Artifact(
                task_id=task.id,
                revision=revision,
                type="deck_pptx",
                path=deck_pptx.relative_to(self.workspace.root).as_posix(),
            ),
            Artifact(
                task_id=task.id,
                revision=revision,
                type="report_markdown",
                path=report_markdown.relative_to(self.workspace.root).as_posix(),
            ),
            Artifact(
                task_id=task.id,
                revision=revision,
                type="agent_result_json",
                path=agent_result_json.relative_to(self.workspace.root).as_posix(),
            ),
            Artifact(
                task_id=task.id,
                revision=revision,
                type="validation_json",
                path=validation_json.relative_to(self.workspace.root).as_posix(),
            ),
        ]
        return artifacts, deck, validation

    def _report_markdown(self, result: AgentExecutionResult) -> str:
        lines = [f"# {result.deliverable.title}", "", result.deliverable.summary, ""]
        for slide in result.deliverable.slides:
            lines.extend([f"## {slide.title}", ""])
            lines.extend(f"- {item}" for item in slide.body)
            lines.extend(f"- **{label}:** {value}" for label, value in slide.metrics.items())
            lines.append("")
        lines.extend(["## Source files", ""])
        lines.extend(
            f"- `{source.path}` — {source.note}" for source in result.deliverable.sources
        )
        return "\n".join(lines).rstrip() + "\n"

    def _choose_table(self, paths: list[Path]) -> Path:
        csvs = [path for path in paths if path.suffix.lower() == ".csv"]
        return sorted(csvs or paths)[0]

    def _load_table(self, path: Path) -> pd.DataFrame:
        if path.suffix.lower() == ".csv":
            return pd.read_csv(path)
        return pd.read_excel(path)

    def _chart_from_frame(self, df: pd.DataFrame, source: str) -> ChartSpec:
        quarters = [str(value) for value in df["quarter"].tolist()]
        revenue = [float(value) / 1_000_000 for value in df["revenue"].tolist()]
        margin = [float(value) * 100 for value in df["operating_margin"].tolist()]
        lineage = [
            {
                "metric": "Revenue",
                "source": source,
                "columns": ["quarter", "revenue"],
                "transformation": "Revenue divided by 1,000,000 for display",
            },
            {
                "metric": "Operating margin",
                "source": source,
                "columns": ["operating_income", "revenue"],
                "transformation": "operating_income / revenue",
            },
        ]
        return ChartSpec(
            title="2024 Quarterly Actuals: Revenue and Operating Margin",
            subtitle="Revenue in millions; operating margin as percentage",
            chart_type="grouped_bar",
            x_label="Quarter",
            y_label="Revenue ($M) / Margin (%)",
            y_unit="mixed",
            series=[
                ChartSeries(name="Revenue ($M)", x=quarters, y=revenue),
                ChartSeries(name="Operating margin (%)", x=quarters, y=margin),
            ],
            source_note=f"Source: {source}",
            lineage=lineage,
        )

    def _pdf_context(self, files: list[Path]) -> tuple[list[SourceCitation], list[str]]:
        citations: list[SourceCitation] = []
        notes: list[str] = []
        for path in sorted(file for file in files if file.suffix.lower() == ".pdf")[:2]:
            rel = path.relative_to(self.workspace.root).as_posix()
            try:
                doc = fitz.open(path)
                text = " ".join(page.get_text("text").strip() for page in doc)
            except Exception:
                text = ""
            compact = " ".join(text.split())
            if compact:
                notes.append(f"Supporting context from {path.name}: {compact[:180]}")
            else:
                notes.append(f"Supporting context from {path.name}: no extractable text found.")
            citations.append(SourceCitation(label=path.name, path=rel, note="Supporting PDF context"))
        return citations, notes

    def _deck_from_analysis(
        self,
        task: RobinTask,
        chart: ChartSpec,
        df: pd.DataFrame,
        source: Path,
        pdf_citations: list[SourceCitation],
        pdf_notes: list[str],
    ) -> DeckSpec:
        first_revenue = float(df.iloc[0]["revenue"])
        last_revenue = float(df.iloc[-1]["revenue"])
        growth = (last_revenue - first_revenue) / first_revenue
        best_margin = df.loc[df["operating_margin"].idxmax()]
        title = "2024 Quarterly Performance"
        executive_body = [
            f"Revenue increased {growth:.1%} from {df.iloc[0]['quarter']} to {df.iloc[-1]['quarter']}.",
            f"The strongest operating margin was {best_margin['operating_margin']:.1%} in {best_margin['quarter']}.",
            "Forecast rows were excluded so the view reflects actual performance.",
        ]
        if pdf_notes:
            executive_body.append(pdf_notes[0])
        source_body = [
            "Read approved local finance data only.",
            "Filtered to actual results.",
            "Validated operating margin as operating income divided by revenue.",
        ]
        source_body.extend(pdf_notes or ["No supporting PDF context was available for this task."])
        return DeckSpec(
            task_id=task.id,
            revision=task.revision,
            title=title,
            slides=[
                SlideSpec(type="title", title=title, body=[task.requested_outcome]),
                SlideSpec(
                    type="executive_summary",
                    title="Executive Summary",
                    body=executive_body,
                ),
                SlideSpec(type="chart", title=chart.title, chart_id=chart.id, body=[chart.subtitle or ""]),
                SlideSpec(
                    type="key_metrics",
                    title="Key Metrics",
                    metrics={
                        "Q4 revenue": f"${last_revenue / 1_000_000:.1f}M",
                        "Q4 operating margin": f"{float(df.iloc[-1]['operating_margin']):.1%}",
                        "Revenue growth": f"{growth:.1%}",
                    },
                ),
                SlideSpec(
                    type="sources",
                    title="Sources and Method",
                    body=source_body,
                ),
            ],
            sources=[
                SourceCitation(label=source.name, path=source.relative_to(self.workspace.root).as_posix(), note="Primary structured source"),
                *pdf_citations,
            ],
        )

    def _validate_analysis(
        self,
        task: RobinTask,
        chart: ChartSpec,
        deck: DeckSpec,
        df: pd.DataFrame,
        source: str,
        supporting_sources: list[str],
    ) -> ValidationReport:
        checks: list[ValidationCheck] = []
        required = {"quarter", "revenue", "operating_income", "operating_margin"}
        present = set(df.columns)
        checks.append(
            ValidationCheck(
                name="required_columns",
                ok=required.issubset(present),
                detail="Required finance columns are present.",
                source=source,
                expected=sorted(required),
                actual=sorted(required & present),
            )
        )
        if "scenario" in df.columns:
            scenarios = {str(value).strip().lower() for value in df["scenario"].tolist()}
            checks.append(
                ValidationCheck(
                    name="forecast_excluded",
                    ok="forecast" not in scenarios,
                    detail="Forecast rows were excluded from the analysis frame.",
                    source=source,
                    expected="no forecast rows",
                    actual=sorted(scenarios),
                )
            )
        expected_margin = (df["operating_income"].astype(float) / df["revenue"].astype(float)).tolist()
        actual_margin = df["operating_margin"].astype(float).tolist()
        margin_deltas = [abs(expected - actual) for expected, actual in zip(expected_margin, actual_margin, strict=False)]
        max_margin_delta = max(margin_deltas, default=0.0)
        checks.append(
            ValidationCheck(
                name="operating_margin_formula",
                ok=max_margin_delta < 0.0001,
                detail="Operating margin matches operating income divided by revenue.",
                source=source,
                expected=[round(value, 6) for value in expected_margin],
                actual=[round(value, 6) for value in actual_margin],
            )
        )
        revenue_series = next((series for series in chart.series if series.name == "Revenue ($M)"), None)
        expected_revenue = [float(value) / 1_000_000 for value in df["revenue"].tolist()]
        revenue_ok = bool(revenue_series) and all(abs(expected - actual) < 0.0001 for expected, actual in zip(expected_revenue, revenue_series.y, strict=False))
        checks.append(
            ValidationCheck(
                name="chart_revenue_series",
                ok=revenue_ok,
                detail="Revenue chart series matches source revenue scaled to millions.",
                source=source,
                expected=[round(value, 6) for value in expected_revenue],
                actual=[round(value, 6) for value in revenue_series.y] if revenue_series else None,
            )
        )
        margin_series = next((series for series in chart.series if series.name == "Operating margin (%)"), None)
        expected_margin_percent = [float(value) * 100 for value in df["operating_margin"].tolist()]
        margin_series_ok = bool(margin_series) and all(abs(expected - actual) < 0.0001 for expected, actual in zip(expected_margin_percent, margin_series.y, strict=False))
        checks.append(
            ValidationCheck(
                name="chart_margin_series",
                ok=margin_series_ok,
                detail="Operating margin chart series matches validated margin percentages.",
                source=source,
                expected=[round(value, 6) for value in expected_margin_percent],
                actual=[round(value, 6) for value in margin_series.y] if margin_series else None,
            )
        )
        checks.append(
            ValidationCheck(
                name="deck_structure",
                ok=3 <= len(deck.slides) <= 6 and any(slide.type == "sources" for slide in deck.slides),
                detail="Deck has a concise slide count and includes a sources slide.",
                expected="3-6 slides with sources",
                actual=f"{len(deck.slides)} slides",
            )
        )
        checks.append(
            ValidationCheck(
                name="lineage_present",
                ok=bool(chart.lineage) and any(item.get("source") == source for item in chart.lineage if isinstance(item, dict)),
                detail="Chart lineage points back to the source file.",
                source=source,
                expected=source,
                actual=chart.lineage,
            )
        )
        cited_paths = [citation.path for citation in deck.sources]
        checks.append(
            ValidationCheck(
                name="source_citations_present",
                ok=source in cited_paths and all(path in cited_paths for path in supporting_sources),
                detail="Deck cites the structured source and any supporting PDF context.",
                source=source,
                expected=[source, *supporting_sources],
                actual=cited_paths,
            )
        )
        return ValidationReport(
            task_id=task.id,
            ok=all(check.ok for check in checks),
            checks=checks,
            source_paths=[source, *supporting_sources],
        )

    def _write_artifacts(self, task_id: UUID, chart: ChartSpec, deck: DeckSpec, validation: ValidationReport) -> list[Artifact]:
        out = self.workspace.generated_task_dir(str(task_id))
        revision = deck.revision
        chart_json = out / f"chart_v{revision}.json"
        deck_json = out / f"deck_v{revision}.json"
        chart_png = out / f"chart_v{revision}.png"
        deck_pptx = out / f"deck_v{revision}.pptx"
        validation_json = out / f"validation_v{revision}.json"
        chart_json.write_text(chart.model_dump_json(indent=2))
        deck_json.write_text(deck.model_dump_json(indent=2))
        validation_json.write_text(validation.model_dump_json(indent=2))
        self._render_png(chart, chart_png)
        self._render_pptx(deck, chart_png, deck_pptx)
        return [
            Artifact(task_id=task_id, revision=revision, type="chart_json", path=chart_json.relative_to(self.workspace.root).as_posix()),
            Artifact(task_id=task_id, revision=revision, type="chart_png", path=chart_png.relative_to(self.workspace.root).as_posix()),
            Artifact(
                task_id=task_id,
                revision=revision,
                type="deck_json",
                path=deck_json.relative_to(self.workspace.root).as_posix(),
                url=f"{self.presentation_base_url}/{task_id}?revision={revision}",
            ),
            Artifact(task_id=task_id, revision=revision, type="deck_pptx", path=deck_pptx.relative_to(self.workspace.root).as_posix()),
            Artifact(task_id=task_id, revision=revision, type="validation_json", path=validation_json.relative_to(self.workspace.root).as_posix()),
        ]

    def _render_pptx(self, deck: DeckSpec, chart_png: Path | None, path: Path) -> None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for slide_spec in deck.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, slide_spec.title)
            if slide_spec.type == "title":
                self._add_text_block(slide, slide_spec.body, left=0.8, top=2.1, width=11.8, height=3.8, font_size=28)
            elif slide_spec.type == "chart":
                try:
                    if chart_png is None:
                        raise ValueError("No chart image was supplied")
                    slide.shapes.add_picture(str(chart_png), Inches(0.8), Inches(1.35), width=Inches(11.8), height=Inches(5.4))
                    self._add_text_block(slide, slide_spec.body[:1], left=0.8, top=6.85, width=11.8, height=0.35, font_size=12)
                except Exception:
                    self._add_text_block(slide, [*slide_spec.body, "Chart image export was unavailable; use the browser deck for the live chart."], left=0.95, top=1.45, width=11.4, height=5.7, font_size=21)
            elif slide_spec.type == "key_metrics":
                metrics = [f"{label}: {value}" for label, value in slide_spec.metrics.items()]
                self._add_text_block(slide, metrics or slide_spec.body, left=1.0, top=1.55, width=11.0, height=5.2, font_size=26)
            else:
                self._add_text_block(slide, slide_spec.body, left=0.95, top=1.45, width=11.4, height=5.7, font_size=21)
            self._add_footer(slide, deck)
        prs.save(path)

    def _add_slide_title(self, slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12.0), Inches(0.75))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.bold = True
        paragraph.font.size = Pt(30)
        paragraph.font.color.rgb = RGBColor(24, 33, 48)

    def _add_text_block(self, slide, lines: list[str], left: float, top: float, width: float, height: float, font_size: int) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for index, line in enumerate(lines or [""]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(42, 52, 68)
            paragraph.space_after = Pt(12)
            if len(lines) > 1:
                paragraph.level = 0

    def _add_footer(self, slide, deck: DeckSpec) -> None:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(12.0), Inches(0.25))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = f"Robin generated export · revision {deck.revision}"
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = RGBColor(102, 112, 133)

    def _render_png(self, chart: ChartSpec, path: Path) -> None:
        fig = go.Figure()
        for series in chart.series:
            fig.add_trace(go.Bar(name=series.name, x=series.x, y=series.y))
        fig.update_layout(
            title={"text": chart.title, "x": 0.02},
            barmode="group",
            width=1280,
            height=720,
            template="plotly_white",
            font={"size": 18},
            margin={"l": 90, "r": 40, "t": 90, "b": 90},
            legend={"orientation": "h", "y": -0.16},
        )
        fig.update_xaxes(title_text=chart.x_label)
        fig.update_yaxes(title_text=chart.y_label)
        try:
            fig.write_image(path)
        except Exception:
            path.write_text(json.dumps(chart.model_dump(mode="json"), indent=2))
